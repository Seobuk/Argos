#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Argos STEP measurement report generator.

Given a STEP/IGES/BREP file, this creates an output folder and a PowerPoint
report (.pptx) with the key measured values: overall dimensions, mass / inertia
(for robot-link dynamics), notable cylinder diameters, and a ready-to-paste
URDF <inertial> block. It drives the Qt-free `argos-cli` (digest / props).

Usage:
  py -3.12 scripts/argos_report.py <file.step> [--density 7850] [--out DIR] [--cli PATH]

Requires: python-pptx  (pip install python-pptx)
"""
import argparse
import json
import os
import subprocess
import sys

KO_FONT = "Malgun Gothic"  # Korean-capable font for PowerPoint rendering

# Offscreen 3D views to render. Each entry:
#   (name, Z-up camera projection vector, Korean label, in-plane dims)
# in-plane dims = (horizontal axis, vertical axis) of the bounding box that the
# orthographic view measures, used to draw dimension annotations on the image.
# None means no dimension annotation (e.g. the isometric view).
# Bounding-box axis -> Korean dimension name (가로/세로/높이).
AXIS_KOR = {"x": "가로(X)", "y": "세로(Y)", "z": "높이(Z)"}

# Views to render. (name, Z-up camera vector, label, in-plane dims (h,v) or None)
# Four isometric corners give a "여러 각도" gallery (turntable-like); front/top/
# side are the orthographic 삼면도 used for the dimensioned drawing. Top is key
# for flat parts (a PCB only "reads" from above) and shows the X-Y footprint.
CAMERA_VIEWS = [
    ("iso",   "1, -1, 1",  "등각 (앞·우·위)", None),
    ("iso2",  "-1, -1, 1", "등각 (앞·좌·위)", None),
    ("iso3",  "-1, 1, 1",  "등각 (뒤·좌·위)", None),
    ("iso4",  "1, 1, 1",   "등각 (뒤·우·위)", None),
    ("front", "0, -1, 0",  "정면도 (앞)",     ("x", "z")),
    ("top",   "0, 0, 1",   "평면도 (위)",     ("x", "y")),
    ("right", "1, 0, 0",   "측면도 (옆)",     ("y", "z")),
]
ISO_VIEWS = ["iso", "iso2", "iso3", "iso4"]


def find_cli(explicit):
    if explicit and os.path.isfile(explicit):
        return os.path.abspath(explicit)
    here = os.path.dirname(os.path.abspath(__file__))
    for c in (
        os.path.join(here, "..", "build", "Release", "argos-cli.exe"),
        os.path.join(here, "..", "dist", "Argos", "argos-cli.exe"),
        "argos-cli.exe",
    ):
        if os.path.isfile(c):
            return os.path.abspath(c)
    return "argos-cli.exe"


def find_conv(explicit=None):
    """Locate mayo-conv.exe (used for headless, offscreen 3D rendering)."""
    if explicit and os.path.isfile(explicit):
        return os.path.abspath(explicit)
    here = os.path.dirname(os.path.abspath(__file__))
    for c in (
        os.path.join(here, "..", "build", "Release", "mayo-conv.exe"),
        os.path.join(here, "..", "dist", "Argos", "mayo-conv.exe"),
    ):
        if os.path.isfile(c):
            return os.path.abspath(c)
    return None


def render_views(conv, step, out_dir, width=1600, height=1200, reuse=False):
    """Render shaded 3D images of the part fully offscreen via mayo-conv.

    mayo-conv writes PPM (the only format OCCT can encode without FreeImage); we
    convert to PNG with Pillow. No GUI window is shown, so this works while the
    machine is in use. Returns {view_name: png_path}; empty if rendering is
    unavailable (the report is still generated, just without images).

    With reuse=True, an existing non-empty view_<name>.png is kept as-is (skips
    mayo-conv) so the PPTX can be rebuilt quickly without re-rendering."""
    images = {}
    todo = []
    for name, cam, _label, _dims in CAMERA_VIEWS:
        png = os.path.join(out_dir, f"view_{name}.png")
        if reuse and os.path.isfile(png) and os.path.getsize(png) > 0:
            images[name] = png
        else:
            todo.append((name, cam))
    if reuse and images:
        print(f"[argos_report] reusing {len(images)} existing render(s)")
    if not todo:
        return images

    if not conv:
        print("[argos_report] mayo-conv not found -> skipping missing 3D images")
        return images
    try:
        from PIL import Image
    except Exception:
        print("[argos_report] Pillow not installed (pip install pillow) -> skipping 3D images")
        return images

    for name, cam in todo:
        ini = os.path.join(out_dir, f"_render_{name}.ini")
        ppm = os.path.join(out_dir, f"_render_{name}.ppm")
        png = os.path.join(out_dir, f"view_{name}.png")
        with open(ini, "w", encoding="utf-8") as f:
            f.write(
                "[export]\n"
                f"Image\\width={width}\n"
                f"Image\\height={height}\n"
                "Image\\backgroundColorStart=#FBFCFE\n"
                "Image\\backgroundColorEnd=#D6E2F0\n"
                "Image\\backgroundGradientFill=Vertical\n"
                f'Image\\cameraOrientation="{cam}"\n'
                "Image\\cameraProjection=Orthographic\n"
                "Image\\GraphicsShapeObjectDriver_displayMode=ShadedWithFaceBoundary\n"
                "Image\\GraphicsMeshObjectDriver_displayMode=Shaded\n"
            )
        try:
            subprocess.run([conv, step, "-e", ppm, "--use-settings", ini, "--no-progress"],
                           capture_output=True, text=True, encoding="utf-8",
                           errors="replace", timeout=900)
        except Exception as e:
            print(f"[argos_report] render {name} failed: {e}")
        if os.path.isfile(ppm) and os.path.getsize(ppm) > 0:
            try:
                Image.open(ppm).save(png)
                images[name] = png
                print(f"[argos_report] rendered {name} -> {os.path.basename(png)}")
            except Exception as e:
                print(f"[argos_report] convert {name} failed: {e}")
        for tmp in (ini, ppm):
            try:
                os.remove(tmp)
            except OSError:
                pass
    return images


def _load_font(size, bold=False):
    from PIL import ImageFont
    for fp in ((r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf"),
               r"C:\Windows\Fonts\malgun.ttf"):
        try:
            return ImageFont.truetype(fp, size)
        except Exception:
            pass
    return ImageFont.load_default()


def _arrow_h(draw, x0, x1, y, color, a=14, w=3):
    draw.line([(x0, y), (x1, y)], fill=color, width=w)
    draw.polygon([(x0, y), (x0 + a, y - a // 2), (x0 + a, y + a // 2)], fill=color)
    draw.polygon([(x1, y), (x1 - a, y - a // 2), (x1 - a, y + a // 2)], fill=color)


def _arrow_v(draw, x, y0, y1, color, a=14, w=3):
    draw.line([(x, y0), (x, y1)], fill=color, width=w)
    draw.polygon([(x, y0), (x - a // 2, y0 + a), (x + a // 2, y0 + a)], fill=color)
    draw.polygon([(x, y1), (x - a // 2, y1 - a), (x + a // 2, y1 - a)], fill=color)


def annotate_image(src, dst, hval, vval, hlabel, vlabel):
    """Draw labelled dimension arrows (horizontal + vertical) onto a rendered
    orthographic view, so the reader can see which direction is width vs height.
    The part silhouette is detected by thresholding against the light background;
    extra padding is added on the left/bottom to hold the dimension lines."""
    try:
        from PIL import Image, ImageDraw
    except Exception:
        return False
    try:
        im = Image.open(src).convert("RGB")
    except Exception:
        return False

    # Part silhouette bbox (background is light: luminance > ~210)
    mask = im.convert("L").point(lambda p: 255 if p < 210 else 0)
    box = mask.getbbox()
    if not box:
        return False

    W, H = im.size
    # Scale all annotation sizes to the image so labels stay legible even when the
    # image is shrunk into a multi-view grid.
    fs = max(30, int(W * 0.040))          # label font size
    aw = max(3, int(W * 0.0045))          # arrow line width
    ah = max(12, int(W * 0.020))          # arrowhead size
    ew = max(1, int(W * 0.0016))          # extension line width
    PL = int(W * 0.20)
    PB = int(W * 0.17)
    PT = int(W * 0.02)
    PR = int(W * 0.03)
    canvas = Image.new("RGB", (W + PL + PR, H + PT + PB), (251, 252, 254))
    canvas.paste(im, (PL, PT))
    x0, y0, x1, y1 = box[0] + PL, box[1] + PT, box[2] + PL, box[3] + PT

    d = ImageDraw.Draw(canvas)
    RED = (193, 57, 43)
    GREY = (120, 124, 128)
    font = _load_font(fs, bold=True)
    gap = int(fs * 0.5)

    def text_w(s):
        try:
            b = d.textbbox((0, 0), s, font=font)
            return b[2] - b[0], b[3] - b[1]
        except Exception:
            return (len(s) * fs // 2, fs)

    # Horizontal dimension (width) along the bottom padding
    if hval is not None:
        ydim = y1 + int(fs * 1.6)
        d.line([(x0, y1 + gap), (x0, ydim + gap)], fill=GREY, width=ew)   # extension lines
        d.line([(x1, y1 + gap), (x1, ydim + gap)], fill=GREY, width=ew)
        _arrow_h(d, x0, x1, ydim, RED, a=ah, w=aw)
        s = f"{hlabel} = {fnum(hval, 1)} mm"
        tw, th = text_w(s)
        tx = max(2, (x0 + x1) // 2 - tw // 2)
        d.rectangle([tx - 8, ydim + gap, tx + tw + 8, ydim + gap + th + 12], fill=(255, 255, 255))
        d.text((tx, ydim + gap + 4), s, fill=RED, font=font)

    # Vertical dimension (height) along the left padding
    if vval is not None:
        xdim = x0 - int(fs * 1.8)
        d.line([(x0 + gap // 2, y0), (xdim - gap, y0)], fill=GREY, width=ew)
        d.line([(x0 + gap // 2, y1), (xdim - gap, y1)], fill=GREY, width=ew)
        _arrow_v(d, xdim, y0, y1, RED, a=ah, w=aw)
        s = f"{vlabel} = {fnum(vval, 1)} mm"
        tw, th = text_w(s)
        lbl = Image.new("RGBA", (tw + 16, th + 16), (255, 255, 255, 235))
        ld = ImageDraw.Draw(lbl)
        ld.text((8, 6), s, fill=RED, font=font)
        lbl = lbl.rotate(90, expand=True)
        ly = max(2, (y0 + y1) // 2 - lbl.height // 2)
        canvas.paste(lbl, (max(2, xdim - gap - lbl.width), ly), lbl)

    try:
        canvas.save(dst)
        return True
    except Exception:
        return False


def annotate_views(images, bbox, out_dir):
    """Produce dimension-annotated copies of the orthographic renders. Returns a
    new images dict where annotatable views point to the *_dim.png versions."""
    if not images or not bbox:
        return images
    dimmap = {name: dims for name, _c, _l, dims in CAMERA_VIEWS}
    result = dict(images)
    for name, raw in images.items():
        dims = dimmap.get(name)
        if not dims:
            continue
        hkey, vkey = dims
        dst = os.path.join(out_dir, f"view_{name}_dim.png")
        if annotate_image(raw, dst, bbox.get(hkey), bbox.get(vkey),
                          AXIS_KOR.get(hkey, hkey.upper()), AXIS_KOR.get(vkey, vkey.upper())):
            result[name] = dst
            print(f"[argos_report] annotated {name} -> {os.path.basename(dst)}")
    return result


def run_cli(cli, args):
    try:
        r = subprocess.run([cli, *args], capture_output=True, text=True,
                           encoding="utf-8", timeout=600)
        return r.stdout or "", r.returncode
    except Exception as e:
        return f"(argos-cli error: {e})", 1


def fnum(v, n=2):
    try:
        return f"{float(v):,.{n}f}"
    except Exception:
        return str(v)


def fsci(v, n=4):
    try:
        return f"{float(v):.{n}e}"
    except Exception:
        return str(v)


def main():
    ap = argparse.ArgumentParser(description="Argos STEP -> PPTX measurement report")
    ap.add_argument("step", help="input CAD file (STEP/IGES/BREP)")
    ap.add_argument("--density", type=float, default=7850.0, help="kg/m^3 (default steel)")
    ap.add_argument("--material", default="", help="material name shown in the report")
    ap.add_argument("--out", default=None, help="output folder (default: <file>_Argos_report)")
    ap.add_argument("--cli", default=None, help="path to argos-cli.exe")
    ap.add_argument("--conv", default=None, help="path to mayo-conv.exe (offscreen renderer)")
    ap.add_argument("--no-images", action="store_true", help="skip 3D rendered images")
    ap.add_argument("--reuse", action="store_true",
                    help="reuse existing digest.json + view_*.png in the output folder "
                         "(fast rebuild of the PPTX without re-measuring/re-rendering)")
    ap.add_argument("--with-mass", action="store_true",
                    help="also include the mass / inertia / URDF slides (humanoid link "
                         "dynamics). Off by default: the report focuses on size + photos.")
    a = ap.parse_args()

    step = os.path.abspath(a.step)
    if not os.path.isfile(step):
        print("[argos_report] file not found:", step)
        sys.exit(1)

    cli = find_cli(a.cli)
    stem = os.path.splitext(os.path.basename(step))[0]
    out_dir = a.out or os.path.join(os.path.dirname(step), stem + "_Argos_report")
    os.makedirs(out_dir, exist_ok=True)

    print(f"[argos_report] cli   = {cli}")
    print(f"[argos_report] file  = {step}")
    print(f"[argos_report] out   = {out_dir}")

    digest_path = os.path.join(out_dir, "digest.json")
    urdf_path = os.path.join(out_dir, "inertial.urdf.xml")
    dg = None
    urdf_txt = ""
    if a.reuse and os.path.isfile(digest_path):
        try:
            dg_txt = open(digest_path, encoding="utf-8").read()
            dg = json.loads(dg_txt)
            if os.path.isfile(urdf_path):
                urdf_txt = open(urdf_path, encoding="utf-8").read()
            print("[argos_report] reusing existing digest.json")
        except Exception:
            dg = None
    if dg is None:
        dg_txt, rc = run_cli(cli, ["digest", step, "--density", str(a.density), "--pretty"])
        try:
            dg = json.loads(dg_txt)
        except Exception as e:
            print("[argos_report] failed to parse digest:", e)
            print(dg_txt[:400])
            sys.exit(1)
        if not dg.get("ok"):
            print("[argos_report] digest failed:", dg.get("error"))
            sys.exit(1)
        with open(digest_path, "w", encoding="utf-8") as f:
            f.write(dg_txt)
        # Mass / inertia / URDF are only computed when explicitly requested.
        if a.with_mass:
            urdf_txt, _ = run_cli(cli, ["props", step, "--density", str(a.density), "--urdf"])
            with open(urdf_path, "w", encoding="utf-8") as f:
                f.write(urdf_txt)

    images = {}
    if not a.no_images:
        images = render_views(find_conv(a.conv), step, out_dir, reuse=a.reuse)
        images = annotate_views(images, dg.get("bboxSize"), out_dir)

    pptx_path = os.path.join(out_dir, stem + "_report.pptx")
    build_pptx(dg, urdf_txt, step, a.density, a.material, pptx_path, images, a.with_mass)
    print(f"[argos_report] PPTX  = {pptx_path}")
    print("[argos_report] done.")


def build_pptx(dg, urdf, step, density, material, path, images=None, with_mass=False):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    images = images or {}
    captions = {n: c for n, _cam, c, _d in CAMERA_VIEWS}

    ACCENT = RGBColor(0x1F, 0x6F, 0xC2)
    INK = RGBColor(0x22, 0x24, 0x26)
    DIM = RGBColor(0x6A, 0x6E, 0x73)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def textbox(slide, l, t, w, h, text, size=14, bold=False, color=INK,
                align=PP_ALIGN.LEFT, font=KO_FONT, mono=False):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Consolas" if mono else font
        return tb

    def title_bar(slide, title, subtitle=""):
        bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        tf = bar.text_frame
        tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = KO_FONT
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = subtitle
            r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0xDD, 0xEA, 0xF8); r2.font.name = KO_FONT

    def add_table(slide, rows, l, t, w, col_fracs, header=True, fontsize=13):
        nrows, ncols = len(rows), len(rows[0])
        gt = slide.shapes.add_table(nrows, ncols, l, t, w, Inches(0.4 * nrows)).table
        for ci, frac in enumerate(col_fracs):
            gt.columns[ci].width = Emu(int(w * frac))
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = gt.cell(ri, ci)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                p = cell.text_frame.paragraphs[0]
                run = p.add_run(); run.text = str(val)
                run.font.size = Pt(fontsize); run.font.name = KO_FONT
                if header and ri == 0:
                    run.font.bold = True; run.font.color.rgb = WHITE
                    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
                else:
                    run.font.color.rgb = INK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF8) if (ri % 2) else WHITE
        return gt

    def place_image(slide, name, l, t, max_w, max_h, caption=None):
        path_img = images.get(name)
        if not path_img or not os.path.isfile(path_img):
            return None
        try:
            from PIL import Image as _Img
            iw, ih = _Img.open(path_img).size
        except Exception:
            iw, ih = 4, 3
        ratio = min(max_w / iw, max_h / ih)
        w, h = int(iw * ratio), int(ih * ratio)
        left = l + (max_w - w) // 2
        pic = slide.shapes.add_picture(path_img, left, t, width=w, height=h)
        pic.line.color.rgb = RGBColor(0xCF, 0xD6, 0xDE)
        pic.line.width = Pt(0.75)
        if caption:
            textbox(slide, l, t + h + Inches(0.03), max_w, Inches(0.35), caption,
                    size=11, color=DIM, align=PP_ALIGN.CENTER)
        return pic

    bb = dg.get("bboxSize") or {}
    counts = dg.get("counts") or {}
    mass = dg.get("mass") or {}
    diam = dg.get("cylinderDiametersMm") or []
    vol_mm3 = mass.get("volume_mm3")
    fname = os.path.basename(step)
    mat_label = material or ("강철(steel)" if abs(density - 7850.0) < 1 else f"{density:g} kg/m³")

    gx, gy, gz = bb.get("x"), bb.get("y"), bb.get("z")

    # ---- Slide 1: title (with the part image) ----
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x15, 0x31, 0x52); bg.line.fill.background()
    textbox(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(1.2),
            "Argos 측정 보고서", size=40, bold=True, color=WHITE)
    sub = f"파일: {fname}"
    if None not in (gx, gy, gz):
        sub += f"\n크기: 가로 {fnum(gx,1)} × 세로 {fnum(gy,1)} × 높이 {fnum(gz,1)} mm"
    textbox(s, Inches(0.85), Inches(2.15), Inches(11.7), Inches(1.2), sub,
            size=16, color=RGBColor(0xCF, 0xDD, 0xEE))
    place_image(s, "iso", Inches(3.4), Inches(3.45), Inches(6.5), Inches(3.7))

    # ---- Slide 2: dimensions (가로/세로/높이) + part image ----
    s = prs.slides.add_slide(blank)
    title_bar(s, "주요 치수 — 가로 · 세로 · 높이", fname)
    rows = [["항목", "값"]]
    if None not in (gx, gy, gz):
        rows.append(["가로 (X)", f"{fnum(gx)} mm"])
        rows.append(["세로 (Y)", f"{fnum(gy)} mm"])
        rows.append(["높이 (Z)", f"{fnum(gz)} mm"])
    if vol_mm3 is not None:
        rows.append(["부피", f"{fnum(float(vol_mm3)/1000.0)} cm³"])
    if mass.get("area_mm2") is not None:
        rows.append(["표면적", f"{fnum(mass.get('area_mm2'))} mm²"])
    rows.append(["솔리드 / 면 / 모서리",
                 f"{counts.get('solids','-')} / {counts.get('faces','-')} / {counts.get('edges','-')}"])
    add_table(s, rows, Inches(0.55), Inches(1.5), Inches(6.4), [0.5, 0.5], fontsize=15)
    place_image(s, "iso", Inches(7.35), Inches(1.55), Inches(5.5), Inches(5.2),
                caption="등각 투상도 (Argos 3D 뷰)")

    # ---- Slide 3: dimensioned drawing (정면·평면·측면 + 등각) ----
    if any(k in images for k in ("front", "top", "right", "iso")):
        s = prs.slides.add_slide(blank)
        title_bar(s, "치수 도면 — 정면 · 평면 · 측면 · 등각",
                  f"{fname}  ·  가로/세로/높이가 이미지에 표시됨")
        grid = [("front", Inches(0.45), Inches(1.30)),
                ("top",   Inches(6.95), Inches(1.30)),
                ("right", Inches(0.45), Inches(4.45)),
                ("iso",   Inches(6.95), Inches(4.45))]
        for name, l, t in grid:
            place_image(s, name, l, t, Inches(5.9), Inches(2.55),
                        caption=captions.get(name, name))

    # ---- Slide 4: multi-angle photo gallery (여러 각도) ----
    iso_present = [v for v in ISO_VIEWS if v in images]
    if len(iso_present) >= 2:
        s = prs.slides.add_slide(blank)
        title_bar(s, "여러 각도 — 등각 뷰", fname)
        cells = [(Inches(0.45), Inches(1.30)), (Inches(6.95), Inches(1.30)),
                 (Inches(0.45), Inches(4.45)), (Inches(6.95), Inches(4.45))]
        for name, (l, t) in zip(iso_present[:4], cells):
            place_image(s, name, l, t, Inches(5.9), Inches(2.55),
                        caption=captions.get(name, name))

    # ---- Slide 5: diameters (largest first; paginated when many) ----
    s = prs.slides.add_slide(blank)
    PER_COL = 14
    MAX_COLS = 2
    shown = diam[:PER_COL * MAX_COLS]
    if len(diam) <= len(shown):
        note = f"총 {len(diam)}개"
    else:
        note = f"총 {len(diam)}개 중 상위 {len(shown)}개 (전체는 digest.json)"
    title_bar(s, "주요 지름 (원통/구멍, 큰 것부터)", f"{fname}  ·  {note}" if diam else fname)
    if diam:
        ncols = max(1, min(MAX_COLS, (len(shown) + PER_COL - 1) // PER_COL))
        col_w = Inches(3.0)
        gap = Inches(0.18)
        for ci in range(ncols):
            chunk = shown[ci * PER_COL:(ci + 1) * PER_COL]
            rows = [["#", "지름(mm)", "반경(mm)"]]
            for k, d in enumerate(chunk):
                rows.append([ci * PER_COL + k + 1, fnum(d, 3), fnum(float(d) / 2.0, 3)])
            add_table(s, rows, Inches(0.5) + ci * (col_w + gap), Inches(1.4), col_w,
                      [0.22, 0.42, 0.36], fontsize=11)
        place_image(s, "top", Inches(7.05), Inches(1.55), Inches(5.8), Inches(5.2),
                    caption="원통/구멍 위치 (평면도)")
    else:
        textbox(s, Inches(0.7), Inches(1.6), Inches(11), Inches(1),
                "원통면이 감지되지 않았습니다.", size=16, color=DIM)

    # ---- Optional: mass / inertia + URDF (humanoid link dynamics) ----
    if with_mass:
        s = prs.slides.add_slide(blank)
        title_bar(s, "질량 특성 (휴머노이드 링크 동역학)", f"{fname}  ·  {mat_label}")
        if mass.get("ok"):
            com = mass.get("com_mm") or {}
            it = mass.get("inertia_com_kg_m2") or {}
            pm = mass.get("principal_moments_kg_m2") or []
            m_rows = [["항목", "값"],
                      ["질량", f"{fnum(mass.get('mass_kg'), 4)} kg  ({fnum(float(mass.get('mass_kg',0))*1000.0)} g)"],
                      ["무게중심 COM (mm)",
                       f"({fnum(com.get('x'),3)}, {fnum(com.get('y'),3)}, {fnum(com.get('z'),3)})"],
                      ["주관성모멘트 (kg·m²)",
                       f"{fsci(pm[0]) if len(pm)>0 else '-'},  {fsci(pm[1]) if len(pm)>1 else '-'},  {fsci(pm[2]) if len(pm)>2 else '-'}"]]
            add_table(s, m_rows, Inches(0.7), Inches(1.35), Inches(11.9), [0.34, 0.66])
            textbox(s, Inches(0.7), Inches(3.55), Inches(6), Inches(0.4),
                    "관성텐서 (COM 기준, kg·m²)", size=14, bold=True, color=ACCENT)
            I = [["", "x", "y", "z"],
                 ["x", fsci(it.get("ixx")), fsci(it.get("ixy")), fsci(it.get("ixz"))],
                 ["y", fsci(it.get("ixy")), fsci(it.get("iyy")), fsci(it.get("iyz"))],
                 ["z", fsci(it.get("ixz")), fsci(it.get("iyz")), fsci(it.get("izz"))]]
            add_table(s, I, Inches(0.7), Inches(4.0), Inches(8.2), [0.16, 0.28, 0.28, 0.28], fontsize=12)
        else:
            textbox(s, Inches(0.7), Inches(1.6), Inches(11), Inches(1),
                    "질량 특성을 계산할 수 없습니다 (솔리드가 아님): " + str(mass.get("error", "")),
                    size=16, color=DIM)

        if urdf and urdf.strip():
            s = prs.slides.add_slide(blank)
            title_bar(s, "URDF <inertial> (로봇 링크)", fname)
            textbox(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.5), urdf.strip(),
                    size=14, mono=True, color=INK)

    prs.save(path)


if __name__ == "__main__":
    main()
